"""Generate Aicraft presentation PowerPoint — updated to match project state."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # dark navy
ACCENT   = RGBColor(0x58, 0xA6, 0xFF)   # blue
GREEN    = RGBColor(0x3F, 0xB9, 0x50)   # pass green
ORANGE   = RGBColor(0xFF, 0xA6, 0x57)   # highlight
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x8B, 0x94, 0x9E)
LIGHT    = RGBColor(0xC9, 0xD1, 0xD9)
DARK_CARD = RGBColor(0x16, 0x1B, 0x22)  # card bg
RED      = RGBColor(0xF8, 0x51, 0x49)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── helpers ──────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    ln = shape.line
    if border:
        ln.color.rgb = border
        ln.width = Pt(1)
    else:
        ln.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.05
    return shape


def add_multiline(slide, left, top, width, height, lines, size=14,
                  color=WHITE, font_name="Consolas", line_space=1.2):
    """Add a text box with multiple lines, each potentially styled differently."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            txt, col, bld, sz = line_info, color, False, size
        else:
            txt = line_info[0]
            col = line_info[1] if len(line_info) > 1 else color
            bld = line_info[2] if len(line_info) > 2 else False
            sz  = line_info[3] if len(line_info) > 3 else size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = font_name
        p.space_after = Pt(2)
        p.space_before = Pt(0)
    return txBox


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(sl, BG)

add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "AICRAFT", size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(sl, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "High-Performance Machine Learning Framework", size=28, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(sl, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
         "Zero dependencies  •  Pure C/C++  •  SIMD-Optimized  •  Edge/Embedded Ready",
         size=16, color=GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")

# Version badge
add_rect(sl, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.5), DARK_CARD, ACCENT)
add_text(sl, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.5),
         "v1.0.0  —  75 tests ✓", size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Architecture
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Architecture", size=36, bold=True, color=ACCENT)

modules = [
    ("platform.h",    "Platform / SIMD detection",          Inches(0.8),  Inches(1.3)),
    ("memory.h",      "Arena & Pool allocators",            Inches(0.8),  Inches(2.2)),
    ("tensor.h",      "N-D Tensors + autograd metadata",    Inches(0.8),  Inches(3.1)),
    ("tensor_ops.h",  "add, sub, mul, div, matmul, reshape",Inches(0.8),  Inches(4.0)),
    ("autograd.h",    "Reverse-mode autodiff (22 ops)",     Inches(0.8),  Inches(4.9)),
    ("simd_math.h",   "SIMD kernels + BLIS GEMM",          Inches(4.8),  Inches(1.3)),
    ("fast_math.h",   "Vectorized exp/sigmoid/tanh",        Inches(4.8),  Inches(2.2)),
    ("layers.h",      "Dense, Conv2D, BN, Dropout, Pool",   Inches(4.8),  Inches(3.1)),
    ("loss.h",        "MSE, CE, BCE",                       Inches(4.8),  Inches(4.0)),
    ("optimizer.h",   "SGD/Adam/AdamW + LR schedulers",     Inches(4.8),  Inches(4.9)),
    ("quantize.h",    "INT8 quantization engine",           Inches(8.8),  Inches(1.3)),
    ("serialize.h",   "Binary model save/load",             Inches(8.8),  Inches(2.2)),
    ("error.h",       "Error codes + callbacks",            Inches(8.8),  Inches(3.1)),
    ("thread_pool.h", "Parallel GEMM threading",            Inches(8.8),  Inches(4.0)),
]

for name, desc, left, top in modules:
    card_w = Inches(3.8)
    card_h = Inches(0.75)
    add_rect(sl, left, top, card_w, card_h, DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
    add_text(sl, left + Inches(0.15), top + Inches(0.05), card_w - Inches(0.3), Inches(0.35),
             name, size=13, bold=True, color=ACCENT, font_name="Consolas")
    add_text(sl, left + Inches(0.15), top + Inches(0.38), card_w - Inches(0.3), Inches(0.3),
             desc, size=10, color=GRAY, font_name="Consolas")

# Footer
add_text(sl, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
         "Header-only  •  14 modules  •  All hot paths inlined  •  Zero function-call overhead",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Test Suite Results
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Test Suite — 75/75 Passed ✓", size=36, bold=True, color=GREEN)

test_sections = [
    ("Tensor Core",                5, GREEN),
    ("SIMD Operations",            7, GREEN),
    ("Matrix Operations",          4, GREEN),
    ("Activations",                4, GREEN),
    ("Layers",                     6, GREEN),
    ("Loss Functions",             3, GREEN),
    ("Autograd",                   6, GREEN),
    ("Optimizers",                 2, GREEN),
    ("Memory Management",          3, GREEN),
    ("Serialization",              1, GREEN),
    ("PRNG",                       3, GREEN),
    ("Integration Tests",          2, GREEN),
    ("Error Handling",             2, GREEN),
    ("Tensor Ops (Sub/Div/Resh.)", 3, GREEN),
    ("Autograd Regression",        3, GREEN),
    ("Gradient Clipping",          2, GREEN),
    ("LR Schedulers",              3, GREEN),
    ("Layer Backward",             3, GREEN),
    ("Conv2D/Pool/BN Bwd",        5, GREEN),
    ("Autograd Sig/Tanh/Soft",    4, GREEN),
    ("INT8 Quantization",          4, GREEN),
]

cols = 3
col_w = Inches(3.8)
row_h = Inches(0.42)
x_start = Inches(0.8)
y_start = Inches(1.2)

for idx, (section, count, col) in enumerate(test_sections):
    c = idx % cols
    r = idx // cols
    x = x_start + c * (col_w + Inches(0.3))
    y = y_start + r * row_h

    add_text(sl, x, y, Inches(2.6), row_h,
             f"✓ {section}", size=11, color=LIGHT, font_name="Consolas")
    add_text(sl, x + Inches(2.8), y, Inches(0.8), row_h,
             f"{count}", size=11, bold=True, color=GREEN, font_name="Consolas",
             align=PP_ALIGN.RIGHT)

# Big result badge
add_rect(sl, Inches(3.5), Inches(5.2), Inches(6.3), Inches(0.9), DARK_CARD, GREEN)
add_text(sl, Inches(3.5), Inches(5.25), Inches(6.3), Inches(0.45),
         "Results: 75 passed, 0 failed, 75 total", size=20, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(sl, Inches(3.5), Inches(5.65), Inches(6.3), Inches(0.4),
         "21 test sections — All passed!", size=16, color=WHITE, align=PP_ALIGN.CENTER, font_name="Consolas")


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — XOR Demo
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "XOR Training Demo — Full Pipeline", size=36, bold=True, color=ACCENT)

# Left: Training loss
add_rect(sl, Inches(0.6), Inches(1.2), Inches(5.5), Inches(5.5), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(0.9), Inches(1.3), Inches(5), Inches(0.4),
         "Training Loss (2→16→16→1)", size=16, bold=True, color=ACCENT)

loss_lines = [
    ("  Epoch        Loss", GRAY, True, 12),
    ("  ─────  ──────────", GRAY, False, 12),
    ("      0    0.280944", LIGHT, False, 12),
    ("    100    0.000329", LIGHT, False, 12),
    ("    200    0.000097", LIGHT, False, 12),
    ("    300    0.000048", LIGHT, False, 12),
    ("    500    0.000019", LIGHT, False, 12),
    ("    700    0.000010", LIGHT, False, 12),
    ("    999    0.000005", GREEN, True, 12),
    ("", WHITE, False, 10),
    ("  Completed in 23.919 ms", ORANGE, True, 13),
    ("  41,808 epochs/sec", ORANGE, False, 13),
]
add_multiline(sl, Inches(0.9), Inches(1.85), Inches(5), Inches(4.5), loss_lines)

# Right: Predictions
add_rect(sl, Inches(6.6), Inches(1.2), Inches(6), Inches(3.5), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(6.9), Inches(1.3), Inches(5), Inches(0.4),
         "Predictions", size=16, bold=True, color=ACCENT)

pred_lines = [
    ("  X1   X2   →   Prediction   Target", GRAY, True, 13),
    ("  ───  ───      ──────────   ──────", GRAY, False, 13),
    ("   0    0   →     0.0034        0    ✓", GREEN, False, 13),
    ("   0    1   →     0.9987        1    ✓", GREEN, False, 13),
    ("   1    0   →     0.9977        1    ✓", GREEN, False, 13),
    ("   1    1   →     0.0015        0    ✓", GREEN, False, 13),
    ("", WHITE, False, 10),
    ("  All predictions correct!", GREEN, True, 14),
]
add_multiline(sl, Inches(6.9), Inches(1.85), Inches(5.5), Inches(3), pred_lines)

# Architecture note
add_rect(sl, Inches(6.6), Inches(5.0), Inches(6), Inches(1.7), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(6.9), Inches(5.1), Inches(5), Inches(0.35),
         "Pipeline Used:", size=14, bold=True, color=ACCENT)
pipeline_lines = [
    ("  Dense(2→16) → ReLU", LIGHT, False, 12),
    ("  Dense(16→16) → ReLU", LIGHT, False, 12),
    ("  Dense(16→1) → Sigmoid", LIGHT, False, 12),
    ("  MSE Loss + Adam optimizer (lr=0.01)", LIGHT, False, 12),
    ("  Arena checkpoint/restore per epoch", ORANGE, False, 12),
]
add_multiline(sl, Inches(6.9), Inches(5.5), Inches(5.5), Inches(1.5), pipeline_lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Benchmarks
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Performance Benchmarks (AVX-512)", size=36, bold=True, color=ACCENT)

benchmarks = [
    ("GEMM 512×512",         "2.112 ms",  "+85.1%", "+86.0%"),
    ("MLP Forward (b=32)",   "0.667 ms",  "+43.5%", "+46.7%"),
    ("Full Train Step",      "4.327 ms",  "+34.9%", "+39.1%"),
    ("4× elem-wise (1M)",   "2.863 ms",  "1465 M elem/s", "—"),
    ("Dot Product (1M)",     "0.339 ms",  "6.18 GFLOPS", "—"),
]

# Table header
y = Inches(1.3)
add_rect(sl, Inches(0.6), y, Inches(12), Inches(0.5), RGBColor(0x1C, 0x22, 0x2B))
headers = [
    (Inches(0.8),  "Benchmark",       Inches(3.2)),
    (Inches(4.2),  "Time",            Inches(2)),
    (Inches(6.4),  "vs PyTorch",      Inches(2.5)),
    (Inches(9.2),  "vs TensorFlow",   Inches(2.5)),
]
for hx, htxt, hw in headers:
    add_text(sl, hx, y, hw, Inches(0.5),
             htxt, size=14, bold=True, color=ACCENT, font_name="Consolas")

for i, (name, time_val, vs_pt, vs_tf) in enumerate(benchmarks):
    y = Inches(1.9 + i * 0.6)
    bg_col = DARK_CARD if i % 2 == 0 else BG
    add_rect(sl, Inches(0.6), y, Inches(12), Inches(0.5), bg_col)
    add_text(sl, Inches(0.8),  y, Inches(3.2), Inches(0.5), name,     size=13, color=WHITE)
    add_text(sl, Inches(4.2),  y, Inches(2),   Inches(0.5), time_val, size=13, color=LIGHT, bold=True)
    add_text(sl, Inches(6.4),  y, Inches(2.5), Inches(0.5), vs_pt,    size=13, color=GREEN, bold=True)
    add_text(sl, Inches(9.2),  y, Inches(2.5), Inches(0.5), vs_tf,    size=13, color=GREEN, bold=True)

# Memory benchmark highlight
add_rect(sl, Inches(0.6), Inches(5.0), Inches(5.8), Inches(1.6), DARK_CARD, ORANGE)
add_text(sl, Inches(0.9), Inches(5.1), Inches(5), Inches(0.4),
         "Memory Allocation (1000× 4KB)", size=16, bold=True, color=ORANGE)
mem_lines = [
    ("Arena Allocator:  0.0319 ms", GREEN, True, 14),
    ("System malloc:    1.0248 ms", RED, False, 14),
    ("", WHITE, False, 8),
    ("Arena speedup:    96.9% faster", ORANGE, True, 16),
]
add_multiline(sl, Inches(0.9), Inches(5.55), Inches(5), Inches(1.2), mem_lines)

# Why faster
add_rect(sl, Inches(6.8), Inches(5.0), Inches(5.8), Inches(1.6), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(7.1), Inches(5.1), Inches(5), Inches(0.4),
         "Why Aicraft is faster:", size=14, bold=True, color=ACCENT)
why_lines = [
    ("• No Python interpreter overhead", LIGHT, False, 11),
    ("• Zero dynamic dispatch / vtable", LIGHT, False, 11),
    ("• Arena allocator (no malloc)", LIGHT, False, 11),
    ("• Hand-tuned SIMD + BLIS GEMM", LIGHT, False, 11),
    ("• Fused ops (FMA, softmax+CE)", LIGHT, False, 11),
]
add_multiline(sl, Inches(7.1), Inches(5.55), Inches(5.5), Inches(1.2), why_lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — SIMD / NEON
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "SIMD — Multi-Architecture Support", size=36, bold=True, color=ACCENT)

archs = [
    ("AVX-512", "16-wide (512-bit)", "6×32 GEMM kernel", "x86_64 server", RGBColor(0x79, 0xC0, 0xFF)),
    ("AVX2",    "8-wide (256-bit)",  "6×16 GEMM kernel", "x86_64 desktop", ACCENT),
    ("SSE",     "4-wide (128-bit)",  "Scalar GEMM",      "Legacy x86",    ORANGE),
    ("NEON",    "4-wide (128-bit)",  "6×8 GEMM kernel",  "ARM / Embedded", GREEN),
    ("Scalar",  "1-wide",            "Scalar fallback",  "Any platform",  GRAY),
]

for i, (arch, width, gemm, target, color) in enumerate(archs):
    y = Inches(1.3 + i * 1.05)
    card_w = Inches(11.8)
    add_rect(sl, Inches(0.8), y, card_w, Inches(0.85), DARK_CARD, color)
    add_text(sl, Inches(1.1), y + Inches(0.08), Inches(2), Inches(0.35),
             arch, size=20, bold=True, color=color, font_name="Consolas")
    add_text(sl, Inches(1.1), y + Inches(0.45), Inches(2), Inches(0.35),
             width, size=11, color=GRAY, font_name="Consolas")
    add_text(sl, Inches(3.5), y + Inches(0.15), Inches(3), Inches(0.5),
             gemm, size=13, color=LIGHT, font_name="Consolas")
    add_text(sl, Inches(7.0), y + Inches(0.15), Inches(3), Inches(0.5),
             target, size=13, color=GRAY, font_name="Consolas")

# NEON detail box
add_rect(sl, Inches(0.8), Inches(6.0), Inches(11.8), Inches(1.0), DARK_CARD, GREEN)
add_text(sl, Inches(1.1), Inches(6.05), Inches(11), Inches(0.45),
         "NEON: Real ARM Intrinsics — Not Just Detection", size=15, bold=True, color=GREEN)
neon_detail = ("vaddq_f32 • vmulq_f32 • vfmaq_f32 • vfmsq_f32 • vrecpeq_f32 • "
               "vmovn_u32 • vaddvq_f32 • vmaxvq_f32 • vcgtq_f32")
add_text(sl, Inches(1.1), Inches(6.5), Inches(11), Inches(0.4),
         neon_detail, size=10, color=GRAY, font_name="Consolas")


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — INT8 Quantization
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "INT8 Quantization — Edge Deployment", size=36, bold=True, color=ACCENT)

# Left: How it works
add_rect(sl, Inches(0.6), Inches(1.2), Inches(6), Inches(3.5), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(0.9), Inches(1.3), Inches(5), Inches(0.4),
         "Quantization Pipeline", size=16, bold=True, color=ACCENT)

quant_lines = [
    ("1. Calibrate", ORANGE, True, 14),
    ("   Scan weights → min/max → scale + zero_point", LIGHT, False, 11),
    ("", WHITE, False, 6),
    ("2. Quantize (float32 → uint8)", ORANGE, True, 14),
    ("   q = clamp(round(x / scale) + zp, 0, 255)", LIGHT, False, 11),
    ("   SIMD-accelerated (NEON + AVX2)", GREEN, False, 11),
    ("", WHITE, False, 6),
    ("3. INT8 Inference", ORANGE, True, 14),
    ("   INT8 matmul → INT32 accumulation", LIGHT, False, 11),
    ("   Dequantize output → float32 + bias", LIGHT, False, 11),
    ("", WHITE, False, 6),
    ("4. Result", ORANGE, True, 14),
    ("   ~4× smaller model, edge-deployable", GREEN, True, 11),
]
add_multiline(sl, Inches(0.9), Inches(1.85), Inches(5.5), Inches(3), quant_lines)

# Right: Size comparison
add_rect(sl, Inches(7.0), Inches(1.2), Inches(5.8), Inches(1.8), DARK_CARD, ORANGE)
add_text(sl, Inches(7.3), Inches(1.3), Inches(5), Inches(0.4),
         "Model Size: 10K Parameters", size=16, bold=True, color=ORANGE)
size_lines = [
    ("  FP32:  40,000 bytes  (39.1 KB)", LIGHT, False, 14),
    ("  INT8:  10,000 bytes  ( 9.8 KB)", GREEN, True, 14),
    ("", WHITE, False, 6),
    ("  Compression: ~4.0×", ORANGE, True, 18),
]
add_multiline(sl, Inches(7.3), Inches(1.85), Inches(5), Inches(1.3), size_lines)

# Right: Test results
add_rect(sl, Inches(7.0), Inches(3.3), Inches(5.8), Inches(1.4), DARK_CARD, GREEN)
add_text(sl, Inches(7.3), Inches(3.4), Inches(5), Inches(0.4),
         "Quantization Tests — 4/4 ✓", size=16, bold=True, color=GREEN)
qtest_lines = [
    ("  ✓ Roundtrip accuracy (max err < 0.02)", GREEN, False, 12),
    ("  ✓ Calibration (scale ≈ 0.00784)", GREEN, False, 12),
    ("  ✓ Quantized dense forward", GREEN, False, 12),
    ("  ✓ Model size estimation (4× compression)", GREEN, False, 12),
]
add_multiline(sl, Inches(7.3), Inches(3.85), Inches(5.5), Inches(1), qtest_lines)

# Bottom: Use cases
add_rect(sl, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.5), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4),
         "Target Platforms", size=16, bold=True, color=ACCENT)
platform_lines = [
    ("  ARM Cortex-M/A  •  Raspberry Pi  •  STM32  •  Android NDK  •  iOS  •  Edge TPU pre-processing",
     LIGHT, False, 13),
    ("  Asymmetric per-tensor affine  •  Zero-point ensures real 0 representation  •  No external runtime",
     GRAY, False, 11),
]
add_multiline(sl, Inches(0.9), Inches(5.8), Inches(11.5), Inches(1), platform_lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — Autograd Deep Dive (NEW)
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Autograd Engine — 22 Backward Ops", size=36, bold=True, color=ACCENT)

# Left: Op list
add_rect(sl, Inches(0.6), Inches(1.2), Inches(5.5), Inches(5.5), DARK_CARD, RGBColor(0x30, 0x36, 0x3D))
add_text(sl, Inches(0.9), Inches(1.3), Inches(5), Inches(0.4),
         "Supported Backward Operations", size=16, bold=True, color=ACCENT)

op_col1 = [
    ("  Arithmetic", ORANGE, True, 13),
    ("    ADD, SUB, MUL, DIV", LIGHT, False, 12),
    ("    SCALE, BIAS_ADD", LIGHT, False, 12),
    ("", WHITE, False, 6),
    ("  Linear Algebra", ORANGE, True, 13),
    ("    MATMUL", LIGHT, False, 12),
    ("", WHITE, False, 6),
    ("  Activations", ORANGE, True, 13),
    ("    RELU, SIGMOID, TANH, SOFTMAX", LIGHT, False, 12),
    ("", WHITE, False, 6),
    ("  Reductions", ORANGE, True, 13),
    ("    SUM, MEAN", LIGHT, False, 12),
    ("", WHITE, False, 6),
    ("  Loss Functions", ORANGE, True, 13),
    ("    MSE_LOSS, CE_LOSS, BCE_LOSS", LIGHT, False, 12),
    ("", WHITE, False, 6),
    ("  Structural / Layers", ORANGE, True, 13),
    ("    FLATTEN, RESHAPE, DROPOUT", LIGHT, False, 12),
    ("    MAXPOOL, BATCHNORM, CONV2D", LIGHT, False, 12),
]
add_multiline(sl, Inches(0.9), Inches(1.85), Inches(5), Inches(5), op_col1)

# Right: Key design points
add_rect(sl, Inches(6.6), Inches(1.2), Inches(6), Inches(2.5), DARK_CARD, ACCENT)
add_text(sl, Inches(6.9), Inches(1.3), Inches(5), Inches(0.4),
         "Design Highlights", size=16, bold=True, color=ACCENT)
design_lines = [
    ("  • Dynamic computational graph", LIGHT, False, 12),
    ("    (no static graph limits)", GRAY, False, 11),
    ("  • Reverse-mode autodiff", LIGHT, False, 12),
    ("    with topological sorting", GRAY, False, 11),
    ("  • Gradient accumulation support", LIGHT, False, 12),
    ("  • In-place gradient clipping", LIGHT, False, 12),
    ("    (L2 norm + value clipping)", GRAY, False, 11),
]
add_multiline(sl, Inches(6.9), Inches(1.85), Inches(5.5), Inches(2), design_lines)

# Right bottom: New backward passes
add_rect(sl, Inches(6.6), Inches(4.0), Inches(6), Inches(2.7), DARK_CARD, GREEN)
add_text(sl, Inches(6.9), Inches(4.1), Inches(5), Inches(0.4),
         "Full Layer Backward Passes", size=16, bold=True, color=GREEN)
layer_bwd_lines = [
    ("  ✓ Conv2D backward", GREEN, False, 12),
    ("    weight, bias, and input grads", GRAY, False, 11),
    ("  ✓ MaxPool2D backward", GREEN, False, 12),
    ("    grad routes only to max elements", GRAY, False, 11),
    ("  ✓ BatchNorm backward", GREEN, False, 12),
    ("    input & gamma grads verified", GRAY, False, 11),
    ("  ✓ Sigmoid / Tanh / Softmax / Mean", GREEN, False, 12),
    ("    numerical gradient checks pass", GRAY, False, 11),
]
add_multiline(sl, Inches(6.9), Inches(4.6), Inches(5.5), Inches(2.2), layer_bwd_lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — Feature Summary / Comparison
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Aicraft vs PyTorch / TensorFlow", size=36, bold=True, color=ACCENT)

comparison = [
    ("Feature",            "PyTorch / TF",               "Aicraft"),
    ("Language",           "Python + C++ backend",       "Pure C/C++"),
    ("Dependencies",       "Hundreds of packages",       "Zero"),
    ("Memory",             "malloc/free per op",         "Arena allocator + checkpoint"),
    ("SIMD",               "Generic BLAS calls",         "Hand-tuned AVX-512/NEON"),
    ("GEMM",               "External MKL/OpenBLAS",      "BLIS-style tiled (built-in)"),
    ("Autograd",           "22 ops (Python overhead)",   "22 ops (zero overhead, inlined)"),
    ("Quantization",       "Separate library (QNNX)",    "Built-in INT8 engine"),
    ("Error Handling",     "Python exceptions",          "Error codes + callbacks"),
    ("Model Size (10K)",   "~40 KB (FP32)",              "~10 KB (INT8)"),
    ("Edge Deployment",    "ONNX export + runtime",      "Native, single header"),
]

for i, (feat, pt, ac) in enumerate(comparison):
    y = Inches(1.2 + i * 0.52)
    is_header = (i == 0)
    bg = RGBColor(0x1C, 0x22, 0x2B) if is_header else (DARK_CARD if i % 2 == 1 else BG)
    add_rect(sl, Inches(0.6), y, Inches(12), Inches(0.46), bg)

    fc = ACCENT if is_header else WHITE
    fc2 = ACCENT if is_header else GRAY
    fc3 = ACCENT if is_header else GREEN
    bld = is_header

    add_text(sl, Inches(0.8),  y, Inches(3),   Inches(0.46), feat, size=13, bold=bld, color=fc)
    add_text(sl, Inches(4.0),  y, Inches(3.8), Inches(0.46), pt,   size=12, color=fc2)
    add_text(sl, Inches(8.2),  y, Inches(4.3), Inches(0.46), ac,   size=12, bold=(not is_header), color=fc3)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Summary / Closing
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG)

add_text(sl, Inches(1), Inches(1.0), Inches(11), Inches(1),
         "AICRAFT", size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

stats = [
    ("14",    "Header-only modules"),
    ("75",    "Tests across 21 sections (100% passing)"),
    ("22",    "Autograd backward ops (fully differentiable)"),
    ("4",     "SIMD architectures (AVX-512/AVX2/SSE/NEON)"),
    ("4×",    "Model compression (INT8 quantization)"),
    ("0",     "External dependencies"),
    ("96.9%", "Arena vs malloc speedup"),
]

for i, (num, desc) in enumerate(stats):
    y = Inches(2.2 + i * 0.58)
    add_text(sl, Inches(3.5), y, Inches(1.5), Inches(0.5),
             num, size=28, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT, font_name="Consolas")
    add_text(sl, Inches(5.3), y, Inches(5), Inches(0.5),
             desc, size=18, color=LIGHT, font_name="Consolas")

add_text(sl, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "Pure C/C++  •  Zero Dependencies  •  Edge/Embedded Ready  •  MIT License",
         size=14, color=GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")


# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════
out_path = r"c:\Users\utente\Downloads\Aicraft\Aicraft_Presentation.pptx"
prs.save(out_path)
print(f"Presentation saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
